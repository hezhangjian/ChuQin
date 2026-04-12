from __future__ import annotations

import tempfile
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


SLIDE_WIDTH_INCHES = 13.333333
SLIDE_HEIGHT_INCHES = 7.5
DEFAULT_RENDER_SCALE = 3.0


@dataclass(slots=True)
class PdfToPptResult:
    output_path: Path
    page_count: int


class PdfToPptConverter:
    def convert(
        self,
        input_pdf: Path,
        output_pptx: Path | None = None,
        work_dir: Path | None = None,
        keep_images: bool = False,
        overwrite: bool = False,
        render_scale: float = DEFAULT_RENDER_SCALE,
    ) -> PdfToPptResult:
        source_path = input_pdf.expanduser().resolve()
        self._validate_source(source_path, render_scale)

        destination_path = self._resolve_output_path(source_path, output_pptx)
        if destination_path.exists() and not overwrite:
            raise FileExistsError(f"Output file already exists: {destination_path}")

        destination_path.parent.mkdir(parents=True, exist_ok=True)

        if work_dir is None:
            with tempfile.TemporaryDirectory(prefix="chuqin-pdf-to-ppt-") as temp_dir:
                temp_path = Path(temp_dir)
                page_images = self._render_pdf_pages(source_path, temp_path, render_scale)
                self._build_pptx(page_images, destination_path)
        else:
            temp_path = work_dir.expanduser().resolve()
            temp_path.mkdir(parents=True, exist_ok=True)
            page_images = self._render_pdf_pages(source_path, temp_path, render_scale)
            self._build_pptx(page_images, destination_path)
            if not keep_images:
                self._cleanup_images(page_images)

        return PdfToPptResult(output_path=destination_path, page_count=len(page_images))

    def _resolve_output_path(self, input_pdf: Path, output_pptx: Path | None) -> Path:
        if output_pptx is None:
            return input_pdf.with_suffix(".pptx")
        return output_pptx.expanduser().resolve()

    def _validate_source(self, input_pdf: Path, render_scale: float) -> None:
        if not input_pdf.exists():
            raise FileNotFoundError(f"PDF file not found: {input_pdf}")
        if input_pdf.suffix.lower() != ".pdf":
            raise ValueError(f"Input file must be a PDF: {input_pdf}")
        if render_scale <= 0:
            raise ValueError("`render_scale` must be greater than 0.")
        self._load_fitz()

    def _render_pdf_pages(self, input_pdf: Path, workspace: Path, render_scale: float) -> list[Path]:
        fitz = self._load_fitz()
        document = fitz.open(str(input_pdf))
        if document.page_count == 0:
            document.close()
            raise ValueError("The input PDF has no pages.")

        page_images: list[Path] = []
        try:
            for index, page in enumerate(document, start=1):
                output_image = workspace / f"page-{index:03d}.png"
                # Render above 1x so fine text survives the later fit-into-slide step more cleanly.
                pixmap = page.get_pixmap(matrix=fitz.Matrix(render_scale, render_scale), alpha=False)
                pixmap.save(str(output_image))
                page_images.append(output_image)
        finally:
            document.close()

        return page_images

    def _build_pptx(self, page_images: list[Path], output_pptx: Path) -> None:
        presentation = Presentation()
        presentation.slide_width = Inches(SLIDE_WIDTH_INCHES)
        presentation.slide_height = Inches(SLIDE_HEIGHT_INCHES)

        blank_layout = presentation.slide_layouts[6]
        slide_width = int(presentation.slide_width)
        slide_height = int(presentation.slide_height)

        for image_path in page_images:
            slide = presentation.slides.add_slide(blank_layout)
            left, top, width, height = self._fit_image(image_path, slide_width, slide_height)
            slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

        presentation.save(str(output_pptx))

    def _load_fitz(self):
        try:
            import fitz
        except ImportError as exc:
            raise RuntimeError(
                "PyMuPDF is required to rasterize PDF pages. Install dependency `pymupdf` to use `chuqin pdf to-ppt`."
            ) from exc
        return fitz

    def _fit_image(self, image_path: Path, slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
        with Image.open(image_path) as image:
            image_width, image_height = image.size

        scale = min(slide_width / image_width, slide_height / image_height)
        fitted_width = int(image_width * scale)
        fitted_height = int(image_height * scale)
        left = int((slide_width - fitted_width) / 2)
        top = int((slide_height - fitted_height) / 2)
        return left, top, fitted_width, fitted_height

    def _cleanup_images(self, page_images: list[Path]) -> None:
        for image_path in page_images:
            image_path.unlink(missing_ok=True)
