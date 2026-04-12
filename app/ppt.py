from __future__ import annotations

import shutil
import tomllib
from dataclasses import dataclass
from dataclasses import field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.config import CONFIG_DIR_NAME
from app.config import CONFIG_ENV_VAR
from app.config import load_config


SLIDE_WIDTH_INCHES = 13.333333
SLIDE_HEIGHT_INCHES = 7.5
DEFAULT_TEMPLATE_DIR_PARTS = ("ppt", "templates", "default")


@dataclass(slots=True)
class TextStyle:
    font_name: str
    font_size: int
    color: str
    bold: bool = False


@dataclass(slots=True)
class TitleSlideLayout:
    title_left: float = 0.9
    title_top: float = 1.55
    title_width: float = 11.5
    title_height: float = 1.3
    subtitle_left: float = 0.9
    subtitle_top: float = 3.05
    subtitle_width: float = 9.5
    subtitle_height: float = 0.7


@dataclass(slots=True)
class PptTemplate:
    title_style: TextStyle
    subtitle_style: TextStyle
    background_color: str = "F7F4ED"
    accent_color: str = "D97706"
    background_image: Path | None = None
    cover_image: Path | None = None
    subtitle_text: str = ""
    layout: TitleSlideLayout = field(default_factory=TitleSlideLayout)


@dataclass(slots=True)
class PptCreateResult:
    output_path: Path
    template_path: Path | None
    used_builtin_template: bool


class PptCreator:
    def create(
        self,
        ppt_name: str,
        output_pptx: Path | None = None,
        template: Path | None = None,
        overwrite: bool = False,
    ) -> PptCreateResult:
        normalized_name = self._normalize_ppt_name(ppt_name)
        destination_path = self._resolve_output_path(normalized_name, output_pptx)
        if destination_path.exists() and not overwrite:
            raise FileExistsError(f"Output file already exists: {destination_path}")

        destination_path.parent.mkdir(parents=True, exist_ok=True)

        template_path = self._resolve_template_path(template)
        if template_path is not None and template_path.suffix.lower() == ".pptx":
            shutil.copyfile(template_path, destination_path)
            return PptCreateResult(
                output_path=destination_path,
                template_path=template_path,
                used_builtin_template=False,
            )

        template_spec, used_builtin_template = self._load_template(template_path)
        self._build_pptx(normalized_name, destination_path, template_spec)
        return PptCreateResult(
            output_path=destination_path,
            template_path=template_path,
            used_builtin_template=used_builtin_template,
        )

    def _normalize_ppt_name(self, ppt_name: str) -> str:
        normalized = ppt_name.strip()
        if not normalized:
            raise ValueError("`ppt_name` cannot be empty.")
        return normalized

    def _resolve_output_path(self, ppt_name: str, output_pptx: Path | None) -> Path:
        if output_pptx is not None:
            return output_pptx.expanduser().resolve()
        return (Path.cwd() / f"{ppt_name}.pptx").resolve()

    def _resolve_template_path(self, template: Path | None) -> Path | None:
        if template is not None:
            return template.expanduser().resolve()

        config = load_config()
        default_template = config.root_path / CONFIG_DIR_NAME
        for part in DEFAULT_TEMPLATE_DIR_PARTS:
            default_template = default_template / part
        if default_template.exists():
            return default_template
        return None

    def _load_template(self, template_path: Path | None) -> tuple[PptTemplate, bool]:
        if template_path is None:
            return self._builtin_template(), True

        if template_path.is_dir():
            return self._load_directory_template(template_path), False

        raise FileNotFoundError(
            f"Template path not found or unsupported: {template_path}. "
            "Use a template directory with `template.toml` or a `.pptx` file."
        )

    def _load_directory_template(self, template_dir: Path) -> PptTemplate:
        template_file = template_dir / "template.toml"
        if not template_file.exists():
            raise FileNotFoundError(f"Template config not found: {template_file}")

        with template_file.open("rb") as file_obj:
            data = tomllib.load(file_obj)

        theme = data.get("theme", {})
        if not isinstance(theme, dict):
            theme = {}
        assets = data.get("assets", {})
        if not isinstance(assets, dict):
            assets = {}
        title = data.get("title", {})
        if not isinstance(title, dict):
            title = {}
        subtitle = data.get("subtitle", {})
        if not isinstance(subtitle, dict):
            subtitle = {}
        cover = data.get("cover", {})
        if not isinstance(cover, dict):
            cover = {}
        layout_data = data.get("layout", {})
        if not isinstance(layout_data, dict):
            layout_data = {}

        return PptTemplate(
            title_style=TextStyle(
                font_name=self._get_str(title, "font_name", "Poppins"),
                font_size=self._get_int(title, "font_size", 28),
                color=self._get_str(title, "color", self._get_str(theme, "text_color", "1F2937")),
                bold=self._get_bool(title, "bold", True),
            ),
            subtitle_style=TextStyle(
                font_name=self._get_str(subtitle, "font_name", "Aptos"),
                font_size=self._get_int(subtitle, "font_size", 14),
                color=self._get_str(subtitle, "color", self._get_str(theme, "muted_text_color", "6B7280")),
                bold=self._get_bool(subtitle, "bold", False),
            ),
            background_color=self._get_str(theme, "background_color", "F7F4ED"),
            accent_color=self._get_str(theme, "accent_color", "D97706"),
            background_image=self._resolve_asset_path(template_dir, assets.get("background_image")),
            cover_image=self._resolve_asset_path(template_dir, assets.get("cover_image")),
            subtitle_text=self._get_str(cover, "subtitle", ""),
            layout=TitleSlideLayout(
                title_left=self._get_float(layout_data, "title_left", 0.9),
                title_top=self._get_float(layout_data, "title_top", 1.55),
                title_width=self._get_float(layout_data, "title_width", 11.5),
                title_height=self._get_float(layout_data, "title_height", 1.3),
                subtitle_left=self._get_float(layout_data, "subtitle_left", 0.9),
                subtitle_top=self._get_float(layout_data, "subtitle_top", 3.05),
                subtitle_width=self._get_float(layout_data, "subtitle_width", 9.5),
                subtitle_height=self._get_float(layout_data, "subtitle_height", 0.7),
            ),
        )

    def _build_pptx(self, ppt_name: str, output_pptx: Path, template: PptTemplate) -> None:
        presentation = Presentation()
        presentation.slide_width = Inches(SLIDE_WIDTH_INCHES)
        presentation.slide_height = Inches(SLIDE_HEIGHT_INCHES)

        blank_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(blank_layout)
        self._apply_background(slide, template, int(presentation.slide_width), int(presentation.slide_height))
        self._add_title_block(slide, ppt_name, template)

        presentation.save(str(output_pptx))

    def _apply_background(self, slide, template: PptTemplate, slide_width: int, slide_height: int) -> None:  # noqa: ANN001
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._parse_color(template.background_color)

        if template.background_image is not None:
            slide.shapes.add_picture(
                str(template.background_image),
                0,
                0,
                width=slide_width,
                height=slide_height,
            )

        if template.cover_image is not None:
            slide.shapes.add_picture(str(template.cover_image), Inches(8.55), Inches(0.8), width=Inches(3.8))

        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.9),
            Inches(0.95),
            Inches(1.4),
            Inches(0.18),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self._parse_color(template.accent_color)
        accent.line.fill.background()

    def _add_title_block(self, slide, ppt_name: str, template: PptTemplate) -> None:  # noqa: ANN001
        title_box = slide.shapes.add_textbox(
            Inches(template.layout.title_left),
            Inches(template.layout.title_top),
            Inches(template.layout.title_width),
            Inches(template.layout.title_height),
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_run = title_paragraph.add_run()
        title_run.text = ppt_name
        self._apply_text_style(title_run, template.title_style)

        subtitle_box = slide.shapes.add_textbox(
            Inches(template.layout.subtitle_left),
            Inches(template.layout.subtitle_top),
            Inches(template.layout.subtitle_width),
            Inches(template.layout.subtitle_height),
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.alignment = PP_ALIGN.LEFT
        subtitle_run = subtitle_paragraph.add_run()
        subtitle_run.text = template.subtitle_text
        self._apply_text_style(subtitle_run, template.subtitle_style)

    def _apply_text_style(self, run, style: TextStyle) -> None:  # noqa: ANN001
        run.font.name = style.font_name
        run.font.size = Pt(style.font_size)
        run.font.bold = style.bold
        run.font.color.rgb = self._parse_color(style.color)

    def _resolve_asset_path(self, template_dir: Path, value: object) -> Path | None:
        if not isinstance(value, str):
            return None
        asset_path = (template_dir / value).resolve()
        if not asset_path.exists():
            raise FileNotFoundError(f"Template asset not found: {asset_path}")
        return asset_path

    def _parse_color(self, value: str) -> RGBColor:
        normalized = value.strip().lstrip("#")
        if len(normalized) != 6:
            raise ValueError(f"Invalid color value: {value}")
        return RGBColor.from_string(normalized.upper())

    def _builtin_template(self) -> PptTemplate:
        default_template_dir = f"{CONFIG_DIR_NAME}/{'/'.join(DEFAULT_TEMPLATE_DIR_PARTS)}"
        return PptTemplate(
            title_style=TextStyle(font_name="Poppins", font_size=28, color="1F2937", bold=True),
            subtitle_style=TextStyle(font_name="Aptos", font_size=14, color="6B7280"),
            background_color="F7F4ED",
            accent_color="D97706",
            subtitle_text=f"Created with chuqin • set `{CONFIG_ENV_VAR}` to use {default_template_dir}",
        )

    def _get_str(self, data: dict[str, object], key: str, default: str) -> str:
        value = data.get(key, default)
        return value if isinstance(value, str) else default

    def _get_int(self, data: dict[str, object], key: str, default: int) -> int:
        value = data.get(key, default)
        return value if isinstance(value, int) else default

    def _get_float(self, data: dict[str, object], key: str, default: float) -> float:
        value = data.get(key, default)
        return float(value) if isinstance(value, (int, float)) else default

    def _get_bool(self, data: dict[str, object], key: str, default: bool) -> bool:
        value = data.get(key, default)
        return value if isinstance(value, bool) else default
