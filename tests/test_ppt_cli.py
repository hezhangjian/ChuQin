import sys
import zipfile
from pathlib import Path

from typer.testing import CliRunner

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from app.cli import app
from app.ppt import DEFAULT_TEMPLATE_DIR_PARTS
from app.ppt import PptCreateResult
from app.ppt import PptCreator


def test_cli_ppt_create_invokes_creator(monkeypatch, tmp_path: Path) -> None:
    runner = CliRunner()
    output_pptx = tmp_path / "roadmap.pptx"
    template_dir = tmp_path / "template"
    template_dir.mkdir()

    def fake_create(self, **kwargs: object) -> PptCreateResult:  # noqa: ANN001
        assert kwargs["ppt_name"] == "Roadmap"
        assert kwargs["output_pptx"] == output_pptx
        assert kwargs["template"] == template_dir
        assert kwargs["overwrite"] is True
        return PptCreateResult(output_path=output_pptx, template_path=template_dir, used_builtin_template=False)

    monkeypatch.setattr("app.cli_ppt.PptCreator.create", fake_create)

    result = runner.invoke(
        app,
        ["ppt", "create", "Roadmap", "-o", str(output_pptx), "--template", str(template_dir), "--overwrite"],
    )

    assert result.exit_code == 0
    assert "PPT created successfully." in result.stdout
    assert str(output_pptx) in result.stdout
    assert str(template_dir) in result.stdout


def test_load_directory_template_reads_theme_and_assets(tmp_path: Path) -> None:
    template_dir = tmp_path / "template"
    template_dir.mkdir()
    background_image = template_dir / "bg.png"
    cover_image = template_dir / "cover.png"
    background_image.write_bytes(b"bg")
    cover_image.write_bytes(b"cover")
    (template_dir / "template.toml").write_text(
        """
[theme]
background_color = "#112233"
accent_color = "445566"
text_color = "FFFFFF"
muted_text_color = "CCCCCC"

[title]
font_name = "DIN"
font_size = 30
bold = true

[subtitle]
font_name = "Aptos"
font_size = 16

[cover]
subtitle = "Q2 planning"

[assets]
background_image = "bg.png"
cover_image = "cover.png"

[layout]
title_left = 1.2
subtitle_top = 3.4
""",
        encoding="utf-8",
    )

    creator = PptCreator()
    template = creator._load_directory_template(template_dir)

    assert template.background_color == "#112233"
    assert template.accent_color == "445566"
    assert template.title_style.font_name == "DIN"
    assert template.title_style.font_size == 30
    assert template.subtitle_text == "Q2 planning"
    assert template.background_image == background_image.resolve()
    assert template.cover_image == cover_image.resolve()
    assert template.layout.title_left == 1.2
    assert template.layout.subtitle_top == 3.4


def test_create_builds_single_slide_pptx_when_default_template_missing(tmp_path: Path, monkeypatch) -> None:
    output_pptx = tmp_path / "demo.pptx"
    creator = PptCreator()

    monkeypatch.setattr("app.ppt.load_config", lambda: type("Config", (), {"root_path": tmp_path})())

    result = creator.create("Demo", output_pptx=output_pptx)

    assert result.output_path == output_pptx
    assert result.template_path is None
    assert result.used_builtin_template is True
    assert output_pptx.exists()

    with zipfile.ZipFile(output_pptx) as archive:
        slide_names = [name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
        assert len(slide_names) == 1


def test_create_copies_pptx_template(tmp_path: Path) -> None:
    from pptx import Presentation

    template_pptx = tmp_path / "template.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(str(template_pptx))

    output_pptx = tmp_path / "copied.pptx"
    creator = PptCreator()

    result = creator.create("Copied", output_pptx=output_pptx, template=template_pptx)

    assert result.template_path == template_pptx.resolve()
    assert result.used_builtin_template is False
    assert output_pptx.read_bytes() == template_pptx.read_bytes()


def test_resolve_template_path_uses_default_templates_directory(tmp_path: Path, monkeypatch) -> None:
    template_dir = tmp_path / ".chuqin"
    for part in DEFAULT_TEMPLATE_DIR_PARTS:
        template_dir = template_dir / part
    template_dir.mkdir(parents=True)

    monkeypatch.setattr("app.ppt.load_config", lambda: type("Config", (), {"root_path": tmp_path})())

    creator = PptCreator()

    assert creator._resolve_template_path(None) == template_dir.resolve()
